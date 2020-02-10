#!/usr/bin/env python3

from pptx import Presentation
import pptx
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import subprocess

icloud_dir="/Users/dillonlue/Library/Mobile Documents/com~apple~CloudDocs/slides"

original_powerpoints = os.listdir(os.path.join(icloud_dir, "originals"))

def create_new_slides(slide_name, icloud_dir):
    slides_dir = os.path.join(icloud_dir, "originals", slide_name)
    prs = Presentation(slides_dir)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape._sp.is_autoshape and shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                if type(shape.fill._fill) != pptx.dml.fill._NoFill:
                    if shape.text == "" or not int(shape.text[0]) not in list(range(0,10)):
                        e = shape.element
                        e.getparent().remove(e)
    prs.save(os.path.join(icloud_dir, "answers", slide_name))
    c = "unoconv -f pdf \"{}\"".format(os.path.join(icloud_dir, "answers", slide_name))
#    print(c)
    subprocess.run([c], shell = True)
    pdf_name = slide_name.replace("pptx", "pdf")
    quiz_dir = os.path.join(icloud_dir, "quizes", pdf_name)
    os.rename(os.path.join(icloud_dir, "answers", pdf_name), quiz_dir)

for powerpoint in original_powerpoints:
    print("Converting {}".format(powerpoint))
    create_new_slides(powerpoint, icloud_dir)
